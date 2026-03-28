import json
from pathlib import Path
from pypdf import PdfReader
from pptx import Presentation

output_dir = Path("output")
output_dir.mkdir(parents=True, exist_ok=True)

# Clean 
def clean_text(raw):
    if not raw: return ""
    return " ".join(raw.replace("\n", " ").split())

# Save file
def save_text(filename,text):
    path = output_dir / filename
    path.write_text(text, encoding="utf-8-sig")
    return str(path)

# Main Program

metadata = []
raw_input = input("Enter file path: ")
path = Path(raw_input)

ext = path.suffix.lower()
doc_name = path.name

# TXT
if ext == ".txt":
        raw_text = path.read_text(encoding="utf-8-sig")
        cleaned = clean_text(raw_text)
    
        metadata.append({
            "source_type": "txt",
            "document": doc_name, 
            "cleaned_text": cleaned
            })
        print(cleaned)
# PDF

elif ext == ".pdf":
    reader = PdfReader(path)
    metadata.append({
            "source_type": "pdf",
             "document": doc_name, 
    })
    for i, page in enumerate(reader.pages, start=1):
        cleaned_page = clean_text(page.extract_text())
        metadata.append({
             "page": i, 
             "cleaned_text": cleaned_page
             })
        print(f"\n--- Page {i} ---\n{cleaned_page}")

# pptx
elif ext == ".pptx":
    pwpt = Presentation(path)
    metadata.append({
            "source_type": "pptx",
            "document": doc_name,
    })
    for i, slide in enumerate(pwpt.slides, start=1):
        text_runs = []

        for shape in slide.shapes:
              if not shape.has_text_frame:
                   continue
              
              for paragraph in shape.text_frame.paragraphs:
                   for run in paragraph.runs:
                        text_runs.append(run.text)

        raw_slide_text = " ".join(text_runs)
        cleaned_text = clean_text(raw_slide_text)
        metadata.append({
            "slide": i,
            "cleaned_text": cleaned_text
        })
        print(f"\n--- Slide {i} ---\n{cleaned_text}")

# Unsupported File Type

else:
     print("Unsupported file type")

# Saving Metadata

meta_path = output_dir / f"metadata_{path.stem}.json" 
meta_path.write_text(json.dumps(metadata, indent=2))

print(f"\nMetadata saved to: {meta_path}")