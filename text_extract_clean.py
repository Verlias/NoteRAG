import json
from pathlib import Path
from pypdf import PdfReader
from pptx import Presentation

# Create a program that can extract at least TXT and PDF files and convert it to raw text
# Clean Text - Remove headers and footers, remove page numbers
# Store extract text locally (file system or database)
# Save basic metadata - Document name, page number (for PDFs)
# Verify correctness by printing clean text to the console
# /Users/zachwang/projectDocument/pdftest.pdf
# /Users/zachwang/projectDocument/txttest.txt
# /Users/zachwang/projectDocument/pptxtest.pptx

# Now add compatability with pptx files

output_dir = Path("output")
output_dir.mkdir(parents=True, exist_ok=True)

# Clean 
def clean_text(raw):
    if not raw: return ""
    return " ".join(raw.replace("\n", " ").split())

# Save file
def save_text(filename,text):
    path = output_dir / filename
    path.write_text(text, encoding="utf-8")
    return str(path)

# Main Program

metadata = []
raw_input = input("Enter file path: ")
path = Path(raw_input)

ext = path.suffix.lower()
doc_name = path.name

# TXT
if ext == ".txt":
        raw_text = path.read_text(encoding="utf-8")
        cleaned = clean_text(raw_text)
        out_path = save_text(f"{doc_name}_cleaned.txt", cleaned)
        metadata.append({"document": doc_name, "output_file": out_path})
        print(cleaned)
# PDF

elif ext == ".pdf":
    reader = PdfReader(path)
    for i, page in enumerate(reader.pages, start=1):
        cleaned_page = clean_text(page.extract_text())
        out_path = save_text(f"{doc_name}_page_{i}.txt", cleaned_page)
        metadata.append({
             "document": doc_name, 
             "page": i, 
             "out_file": out_path
             })
        print(f"\n--- Page {i} ---\n{cleaned_page}")

# pptx
elif ext == ".pptx":
    pwpt = Presentation(path)

    for i, slide in enumerate(pwpt.slides, start=1):
        text_runs = []

        for shape in slide.shapes:
              if not shape.has_text_frame:
                   continue
              
              for paragraph in shape.text_frame.paragraphs:
                   for run in paragraph.runs:
                        text_runs.append(run.text)

        raw_slide_text = " ".join(text_runs)
        cleaned_text = clean_text(raw_slide_text)
        out_path = save_text(f"{path.stem}_slide_{i}.txt", cleaned_text)
        metadata.append({
            "document": doc_name,
            "slide": i,
            "out_file": out_path
        })
        print(f"\n--- Slide {i} ---\n{cleaned_text}")

# Unsupported File Type

else:
     print("Unsupported file type")

# Saving Metadata

meta_path = output_dir / f"metadata_{path.stem}.json" # path.stem is name without extension
meta_path.write_text(json.dumps(metadata, indent=2))

print(f"\nMetadata saved to: {meta_path}")